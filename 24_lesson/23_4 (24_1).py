from random import randint, choice, random, shuffle, randrange
from pptx.util import Pt

name = ['randint', 'choice', 'random', 'shuffle', 'randrange']
x = ['randint(a, b) method of random.Random instance',
     'choice(seq) method of random.Random instance Choose a random element from a non-empty sequence.',
     'random() method of random.Random instance random() -> x in the interval [0, 1).',
     'Shuffle list x in place, and return None.Optional argument random is a 0-argument function returning a '
     'random float in [0.0, 1.0); if it is the default None, the '
     'standard random.random will be used. ',
     'randrange(start, stop=None, step=1) method of random.Random instance Choose a random item from '
     'range(start, stop[, step]).']

from pptx import Presentation

prs = Presentation()  # создали презентацию
title_slide_layout = prs.slide_layouts[1]
for i in range(5):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.placeholders[0]
    subtitle = slide.placeholders[1]
    title.text = name[i]
    subtitle.text = x[i]
    title.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.name = 'Courier New'
    title.text_frame.paragraphs[0].font.name = 'Courier New'

prs.save('test11.pptx')